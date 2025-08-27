#!/usr/bin/env python
# coding: utf-8

# In[23]:


# -----------------------------------------------
# Script para generar presentación mensual en PPT
# Requiere instalar:
#   !pip install python-pptx pandas
# -----------------------------------------------

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import pandas as pd
from datetime import date
import calendar
import locale

# Asegúrate de usar español para nombres de meses
locale.setlocale(locale.LC_TIME, "es_MX.UTF-8")

# Obtener el último mes completo
hoy = date.today()
mes_previo = hoy.month - 2 
ultimo_mes = hoy.month - 1 
anio_actual = hoy.year

# Si estamos en enero, restar un año y usar diciembre del año anterior
if ultimo_mes == 0:
    ultimo_mes = 12
    anio_actual -= 1

# Formatear nombre del mes en español
nombre_mes = calendar.month_name[ultimo_mes].capitalize()
texto_periodo = f"Enero - {nombre_mes}"

nombre_mes_previo = calendar.month_name[mes_previo].capitalize()


# Cargar plantilla existente
prs = Presentation("../auxiliares/Plantilla1.pptx")

# Leer el CSV con títulos
df_titulos = pd.read_csv("../auxiliares/titulos_tarjetas.csv")

# Ruta a imágenes
ruta_imagenes = "../salidas"

# Mapeo diapos → imágenes
mapeo = {
    1: ["tabla_variacion_mensual_absolutos.png"],
    2: ["tabla_variacion_mensual_promedios.png"],
    3: ["Tbl_acumulados.png"],
    4: ["TblAcum_AI.png", "TblMes_AI.png", "VarAcum_2019T_AI.png", "VarMes_2019T_AI.png"],
    5: ["grafico_acum_AI.png", "VarAcum_2019G_AI.png", "VarAcum_2024G_AI.png"],
    6: ["grafico_mes_AI.png", "VarMes_2019G_AI.png", "VarMes_2024G_AI.png"],
    7: ["graficas_combinadas.png"],
    8: ["grafico_alcaldias.png"],
    9: ["grafico_sectores.png"],
    10: ["grafico_cuadrantes.png"],
    11: ["TblAcum_HD.png", "TblMes_HD.png", "VarAcum_2019T_HD.png", "VarMes_2019T_HD.png"],
    12: ["grafico_acum_HD.png", "VarAcum_2019G_HD.png", "VarAcum_2024G_HD.png"],
    13: ["grafico_mes_HD.png", "VarMes_2019G_HD.png", "VarMes_2024G_HD.png"],
    14: ["grafico_acum_HDV.png", "VarAcum_2019G_HDV.png", "VarAcum_2024G_HDV.png"],
    15: ["grafico_mes_HDV.png", "VarMes_2019G_HDV.png", "VarMes_2024G_HDV.png"],
    16: ["TblAcum_LD.png", "TblMes_LD.png", "VarAcum_2019T_LD.png", "VarMes_2019T_LD.png"],
    17: ["grafico_acum_LD.png", "VarAcum_2019G_LD.png", "VarAcum_2024G_LD.png"],
    18: ["grafico_mes_LD.png", "VarMes_2019G_LD.png", "VarMes_2024G_LD.png"],
    19: ["barras_homi_les_alc.png"],
    20: ["TblAcum_RV.png", "TblMes_RV.png", "VarAcum_2019T_RV.png", "VarMes_2019T_RV.png"],
    21: ["grafico_acum_RV.png", "VarAcum_2019G_RV.png", "VarAcum_2024G_RV.png"],
    22: ["grafico_mes_RV.png", "VarMes_2019G_RV.png", "VarMes_2024G_RV.png"],
    23: ["mapa_interactivo_homicidio_lesiones.png"]
}

# Diccionario que asigna slide_index → ID del delito (mismo que en el CSV)
titulos_por_slide = {
    4: "AI",   # diapo 5
    5: "AI",   # ← diapo de gráfica para AI
    6: "AI",   # ← diapo de gráfica mes para AI
    11: "HD",  # diapo 12
    12: "HD",  # ← diapo de gráfica para HD
    13: "HD",  # ← diapo de gráfica mes para HD
    14: "HDV",  # ← diapo de gráfica para HDV
    15: "HDV",  # ← diapo de gráfica mes para HDV
    16: "LD",  # diapo 17
    17: "LD",  # ← diapo de gráfica para LD
    18: "LD",  # ← diapo de gráfica mes para LD
    20: "RV",  #diapo 21
    21: "RV",  # ← diapo de gráfica para RV
    22: "RV"  # ← diapo de gráfica para RV
}

# Función para insertar una imagen en el centro
def insertar_una_imagen(slide, imagen_path):
    slide.shapes.add_picture(imagen_path, Inches(0.5), Inches(1.7), width=Inches(12.5))

# Función para insertar las 4 imágenes en layout tipo tabla + tarjeta
def insertar_tabla_tarjeta(slide, imagenes_paths):
    # Asignar nombres para claridad
    tabla_semestral, tabla_mensual, tarjeta_semestral, tarjeta_mensual = imagenes_paths

    # Posiciones pensadas para una slide de 13.33 x 7.5 pulgadas
    slide.shapes.add_picture(tabla_semestral, left=Inches(0.5), top=Inches(4.2), width=Inches(9.5))
    slide.shapes.add_picture(tabla_mensual,   left=Inches(0.5), top=Inches(6), width=Inches(9.5))
    slide.shapes.add_picture(tarjeta_semestral, left=Inches(10), top=Inches(1.5), width=Inches(3.5))
    slide.shapes.add_picture(tarjeta_mensual,   left=Inches(10), top=Inches(4.3), width=Inches(3.5))

def insertar_grafica_tarjetas(slide, imagenes_paths):
    grafica, tarjeta_2019, tarjeta_2024 = imagenes_paths

    # Insertar gráfica centrada
    slide.shapes.add_picture(grafica, left=Inches(0.5), top=Inches(2.5), width=Inches(9.5))

    # Insertar tarjetas a la derecha
    slide.shapes.add_picture(tarjeta_2019, left=Inches(10.0), top=Inches(2.0), width=Inches(3.5))
    slide.shapes.add_picture(tarjeta_2024, left=Inches(10.0), top=Inches(4.0), width=Inches(3.5))

# ───────────────────────────────────────────────
# Reemplazar "xxx" en la diapositiva 1 (slide 0)
# ───────────────────────────────────────────────
slide_portada = prs.slides[0]
for shape in slide_portada.shapes:
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "XXX" in run.text:
                    run.text = run.text.replace("XXX", nombre_mes) 

# ───────────────────────────────────────────────
# Reemplazar "bbb" en las diapositivas 2, 3, 4, 9 y última (índices 1, 2, 3)
# ───────────────────────────────────────────────
for i in [1, 2, 3, 8, 23]:  # índices de las diapos
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "mmm" in run.text:
                        run.text = run.text.replace("mmm", nombre_mes)
                    if "aaa" in run.text:
                        run.text = run.text.replace("aaa", nombre_mes_previo)


# ───────────────────────────────────────────────
# Insertar imágenes (gráficos, tablas y tarjetas) según mapeo
# ───────────────────────────────────────────────
for slide_index, imagenes in mapeo.items():
    slide = prs.slides[slide_index]  # Ojo: slide_index empieza en 0
    imagenes_paths = [os.path.join(ruta_imagenes, img) for img in imagenes]

    if len(imagenes) == 1:
        insertar_una_imagen(slide, imagenes_paths[0])

    elif len(imagenes) == 3:
        insertar_grafica_tarjetas(slide, imagenes_paths)

        if slide_index in titulos_por_slide:
            print("▶ Procesando slide (gráfica):", slide_index)
            id_delito = titulos_por_slide[slide_index]
            filtro = df_titulos[df_titulos["archivo"] == id_delito]

            if not filtro.empty:
                fila = filtro.iloc[0]

                # Solo si el título acumulado y mensual son iguales
                #if fila["titulo_acumulado"] == fila["titulo_mensual"]:
                # Reemplazo de "xxx" por el texto del periodo
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "xxx" in run.text:
                                    run.text = run.text.replace("xxx", texto_periodo)
                                if "mmm" in run.text:
                                    run.text = run.text.replace("mmm", nombre_mes)

                # Insertar flecha si existe
                if "flecha" in fila and pd.notna(fila["flecha"]):
                    ruta_flecha = fila["flecha"]
                    slide.shapes.add_picture(ruta_flecha, left=Inches(9.5), top=Inches(1.7), width=Inches(1.5))


    elif len(imagenes) == 4:
        insertar_tabla_tarjeta(slide, imagenes_paths)

        if slide_index in titulos_por_slide:
            print("▶ Procesando slide:", slide_index)
            id_delito = titulos_por_slide[slide_index]
            filtro = df_titulos[df_titulos["archivo"] == id_delito]

            if not filtro.empty:
                fila = filtro.iloc[0]

                #if fila["titulo_acumulado"] == fila["titulo_mensual"]:
                texto = fila["titulo_mensual"]

                # Buscar caja de texto que contenga "xxxxx" y reemplazarlo
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "xxxxx" in run.text:
                                    run.text = run.text.replace("xxxxx",texto)
                                    # Insertar flecha si existe en el CSV
                                    if "flecha" in fila and pd.notna(fila["flecha_mensual"]):
                                        ruta_flecha = fila["flecha_mensual"]
                                        slide.shapes.add_picture(ruta_flecha, left=Inches(8), top=Inches(2), width=Inches(2))

                # Reemplazo adicional de "xxxx" por el texto del periodo
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "xxxx" in run.text:
                                    run.text = run.text.replace("xxxx", texto_periodo)
    else:
        raise ValueError(f"Slide {slide_index+1} tiene {len(imagenes)} imágenes, y no tengo una función para eso aún.")


# Guardar presentación final
prs.save("../salidas/presentacion_mensual.pptx")


# In[ ]:




